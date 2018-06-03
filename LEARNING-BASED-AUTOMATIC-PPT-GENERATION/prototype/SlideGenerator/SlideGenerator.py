import operator
from pptx import Presentation
from pptx.util import Inches, Pt
from nltk import word_tokenize
import re
from textblob import TextBlob
from gensim.summarization import keywords
from nltk.stem import WordNetLemmatizer
from collections import Counter, OrderedDict
import glob
import os


class SlideGenerator:
    def __init__(self):
        self.prs = Presentation('st.pptx')

    def lemmatize_word(self, word, pos='n'):
        wordnet_lemmatizer = WordNetLemmatizer()
        return wordnet_lemmatizer.lemmatize(word, pos=pos)

    def get_keywords(self, text, ratio=0.1):
        return keywords(text, ratio=ratio, split=True)

    def get_nouns(self, text):
        blob = TextBlob(text)
        return list(blob.noun_phrases)

    def get_bullet_title(self, sentences):
        text = ' '.join(sentences)
        keywords = self.get_keywords(text, ratio=0.1)
        if keywords:
            return keywords[0]
        return "<Please Fill in an appropriate title>"

    def get_cleaned_bullets(self, sentences):
        c_sentences = ['']
        for sentence in sentences:
            words = word_tokenize(sentence)
            sentence = re.sub(r'^\W+', '', sentence)
            sentence = re.sub(r'^[0-9\.,?:;!\)\}\]]*', '', sentence) #starting with numbers, symbols, etc is removed
            sentence = re.sub(r'[0-9]*$', '', sentence)
            sentence = re.sub(r'\W+$', '', sentence)
            sentence = re.sub(r'\W+$', '', sentence)

            sentence = re.sub(r'- ', '', sentence)
            sentence = re.sub(r'^[\w+ \t\n,:;?]*[\)\}\]]', '', sentence) #no starting braces in the start
            sentence = re.sub(r'[\(\{\[][\w+ \t\n,:;?]*$', '', sentence)#no ending paranthesis in the end

            sentence = sentence.strip()
            if len(sentence) > 1:
                sentence = sentence[0].upper() + sentence[1:]
            c_sentences.append(sentence)
        return c_sentences

    def create_title_slide(self, inp_title="", inp_subtitle=""):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = inp_title
        subtitle.text = inp_subtitle

    def set_footer(self, slide_no, footer_text):
        slide = self.prs.slides[slide_no]
        left = Inches(2.8)
        width = Inches(3.7)
        top = Inches(6.6)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box_frame = text_box.text_frame
        footer_paragraph = text_box_frame.add_paragraph()
        footer_paragraph.text = footer_text
        footer_paragraph.font.bold = True
        footer_paragraph.font.size = Pt(13)

    def add_bullet_slide(self, title, array_of_bullet_sentences):
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = title
        tf = body_shape.text_frame
        if len(array_of_bullet_sentences) > 0:
            tf.text = array_of_bullet_sentences[0]
            for bullet in array_of_bullet_sentences[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

    def add_text_slide(self, text_array, title=""):
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # left = top = width = height = Inches(1)
        left = Inches(0.5)
        width = Inches(5)
        top = Inches(1)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(35)

        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(20)

        for sentence in text_array:
            p = tf.add_paragraph()
            p.text = sentence
            p.font.size = Pt(20)

    def add_image(self, image_path, inp_title=""):
	bullet_slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
	#title = slide.shapes.title
        shapes = slide.shapes
        title_shape = shapes.title
        #body_shape = shapes.placeholders[1]
        title_shape.text = inp_title
	

	left = Inches(3)
	top = Inches(2)
	height = Inches(4.5)
	#shape = shapes.add_shape(image_path, left, top, 4, height)
        slide.shapes.add_picture(image_path, left, top,height)

    def set_logo(self, slide_no, logo_path):
        if logo_path != '':
            slide = self.prs.slides[slide_no]
            left = Inches(0.2)
            top = Inches(0.2)
            slide.shapes.add_picture(logo_path, left, top)

    def create_presentation(self, output_file, title='', sub_title='Processes', footer='AISSMS COE, PUNE', logo='logo.png', contents=[]):
	sub_title = sub_title.replace("\\","")
        self.create_title_slide(title, sub_title)
	sno=1
	ibullets=[]
	cont_len=len(contents)
        for i in range(0, cont_len, 5):
            bullets = contents[i:i + 5]
	    #img_count=0
	    
			
    	    #sno=sno+1
	    #print sno
            bullets.insert(0, "")
            bullet_title = self.get_bullet_title(bullets)
            if bullet_title:
                bullet_title = bullet_title.title()
            bullets = self.get_cleaned_bullets(bullets)
            self.add_bullet_slide(bullet_title, bullets)
            self.set_logo(sno, logo)
            self.set_footer(sno, footer)
	    sno=sno+1
	    for bullet in bullets:
		figures = re.findall(r'([f,F]igure \d{0,9}[_]\d{0,9} [_]*[a-z,A-Z,0-9]+[-,_]?[a-zA-z,0-9]+)', bullet)
		for fig in figures:
    			#print fig,"on slide no. : ", str(sno)
			#img_count=img_count+1
			#ibullets.append(fig.replace('f','F')+" referred on slide no. : "+ str(sno))
			fig= fig.split(' ')[len(fig.split(' '))-1]+'-'+fig.split(' ')[len(fig.split(' '))-2].replace(':','_')
			print fig,' at : ','../src/'+fig+".??g"
			for fig in sorted(glob.glob('uploads/slides-archive/'+fig+".??g")):
    				print "figures added : ",fig
                                 			
    				#print "figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].replace("_",":")
				self.add_image(fig,"Figure "+ fig.split('-')[len(fig.split('-'))-1].split('.')[0])
				print os.getcwd()
				cmd = 'rm -f '+fig
				print cmd
				os.system(cmd)
				#fig_no=fig_no+1
        			self.set_logo(sno, logo)
        			self.set_footer(sno, footer)
				sno = sno + 1 
	    #print bullets[0]
	c_sno=sno
	ibullets.append("All relevant non-referred useful images are shown in upcoming slides.")
	for fc in range(0, len(ibullets), 10):
		ibullet= ibullets[fc:fc+10]
		ibullet.insert(0, "")
		ibullet.insert(1, "")
		ibullet.insert(2, "")
		ibullet.insert(3, "")
		ibullet.insert(4, "")
		#ibullet = self.get_cleaned_bullets(ibullet)
		self.add_bullet_slide('Referred Image(s)', ibullet)
		self.set_logo(c_sno, logo)
            	self.set_footer(c_sno, footer)
		c_sno=c_sno+1		


        fig_slide = c_sno
	fno=1
	for file_name in sorted(glob.glob('uploads/slides-archive/*.pdf'), key=os.path.getatime):
		fname = file_name.split('.')[0]
		'''for img in sorted(glob.glob(fname+"-?*_?*.??g"), key=os.path.getatime):
    				#print img
   			
    				#print "figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].replace("_",":")
				self.add_image(img,"Figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].replace("_",":")+" (f_no."+str(fno)+")")
				#fig_no=fig_no+1
        			self.set_logo(fig_slide, logo)
        			self.set_footer(fig_slide, footer)
				fig_slide=fig_slide+1 
		fno=fno+1'''
		if("tutorial" in fname):
			for img in sorted(glob.glob(fname+"-?*_?*.??g"), key=os.path.getatime):
    				#print img
   			       
    				#print "figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].replace("_",":")
				if(img.split('-')[len(img.split('-'))-1].split('.')[0].split('_')[1] != '1'):
					imno= int(img.split('-')[len(img.split('-'))-1].split('.')[0].split('_')[1])
					#imno=imno-1			
					self.add_image(img,"Figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].split('_')[0]+"_"+str(imno)+" (f_no."+str(fno)+")")
					#fig_no=fig_no+1
        				self.set_logo(fig_slide, logo)
        				self.set_footer(fig_slide, footer)
					fig_slide=fig_slide+1 
			fno=fno+1
		else:
			for img in sorted(glob.glob(fname+"-?*_?*.??g"), key=os.path.getatime):
    				#print img
   			
    				#print "figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0].replace("_",":")
				self.add_image(img,"Figure "+ img.split('-')[len(img.split('-'))-1].split('.')[0]+" (f_no."+str(fno)+")")
				#fig_no=fig_no+1
        			self.set_logo(fig_slide, logo)
        			self.set_footer(fig_slide, footer)
				fig_slide=fig_slide+1 
			fno=fno+1

        self.prs.save(output_file + '.pptx')
